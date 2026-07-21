from pathlib import Path
from pypdf import PdfReader
from rag.cleaner import clean_text


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xls", ".txt", ".csv", ".pptx"}


# ==========================
# PDF
# ==========================

def load_pdf(file_path: str):
    reader = PdfReader(file_path)
    source = Path(file_path).name
    pages = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if not text:
            continue
        pages.append({
            "source": source,
            "page": i + 1,
            "text": clean_text(text)
        })

    return pages


# ==========================
# Word (.docx)
# ==========================

def load_docx(file_path: str):
    from docx import Document

    doc = Document(file_path)
    source = Path(file_path).name
    pages = []

    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        pages.append({
            "source": source,
            "page": i + 1,
            "text": clean_text(text)
        })

    return pages


# ==========================
# Excel (.xlsx / .xls)
# ==========================

def load_excel(file_path: str):
    import pandas as pd

    source = Path(file_path).name
    pages = []
    xl = pd.ExcelFile(file_path)

    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name)
        text = df.to_string(index=False)
        if not text.strip():
            continue
        pages.append({
            "source": source,
            "page": sheet_name,
            "text": clean_text(text)
        })

    return pages


# ==========================
# Plain Text (.txt)
# ==========================

def load_txt(file_path: str):
    source = Path(file_path).name

    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()

    if not text.strip():
        return []

    return [{
        "source": source,
        "page": 1,
        "text": clean_text(text)
    }]


# ==========================
# CSV (.csv)
# ==========================

def load_csv(file_path: str):
    import pandas as pd

    source = Path(file_path).name
    df = pd.read_csv(file_path)
    text = df.to_string(index=False)

    if not text.strip():
        return []

    return [{
        "source": source,
        "page": 1,
        "text": clean_text(text)
    }]


# ==========================
# PowerPoint (.pptx)
# ==========================

def load_pptx(file_path: str):
    from pptx import Presentation

    prs = Presentation(file_path)
    source = Path(file_path).name
    pages = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        text = "\n".join(texts)
        if not text:
            continue
        pages.append({
            "source": source,
            "page": i + 1,
            "text": clean_text(text)
        })

    return pages


# ==========================
# Dispatcher
# ==========================

def load_document(file_path: str):
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext == ".docx":
        return load_docx(file_path)
    elif ext in (".xlsx", ".xls"):
        return load_excel(file_path)
    elif ext == ".txt":
        return load_txt(file_path)
    elif ext == ".csv":
        return load_csv(file_path)
    elif ext == ".pptx":
        return load_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
